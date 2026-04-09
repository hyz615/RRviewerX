from typing import List, Optional
import base64
import io
import mimetypes
import re

import pandas as pd
from bs4 import BeautifulSoup
from docx import Document
from pypdf import PdfReader

from ..core.config import settings
from .agent_service import summarize_visual_material

try:
    import fitz  # type: ignore
except Exception:
    fitz = None

try:
    import os
    import pytesseract
    from PIL import Image

    _tc = os.getenv("TESSERACT_CMD")
    if _tc:
        try:
            pytesseract.pytesseract.tesseract_cmd = _tc
        except Exception:
            pass
except Exception:
    pytesseract = None
    Image = None


def read_pdf_bytes(data: bytes) -> str:
    try:
        reader = PdfReader(io.BytesIO(data))
        texts = []
        for page in reader.pages:
            texts.append(page.extract_text() or "")
        return "\n".join(texts).strip()
    except Exception:
        return ""


def read_docx_bytes(data: bytes) -> str:
    try:
        doc = Document(io.BytesIO(data))
        return "\n".join(p.text for p in doc.paragraphs)
    except Exception:
        return ""


def read_xlsx_bytes(data: bytes) -> str:
    try:
        with pd.ExcelFile(io.BytesIO(data)) as workbook:
            texts = []
            for sheet in workbook.sheet_names:
                texts.append(workbook.parse(sheet).to_csv(index=False))
        return "\n".join(texts)
    except Exception:
        return ""


def read_pptx_bytes(data: bytes) -> str:
    try:
        from pptx import Presentation

        presentation = Presentation(io.BytesIO(data))
        chunks: List[str] = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    chunks.append(str(shape.text))
                try:
                    if getattr(shape, "has_table", False):
                        table = shape.table
                        for row in table.rows:
                            cells = []
                            for cell in row.cells:
                                try:
                                    cells.append(cell.text or "")
                                except Exception:
                                    cells.append("")
                            if any(text.strip() for text in cells):
                                chunks.append("\t".join(cells))
                except Exception:
                    pass
                try:
                    if getattr(shape, "shape_type", None) and str(shape.shape_type).lower().endswith("picture"):
                        image = getattr(getattr(shape, "image", None), "blob", None)
                        if image and pytesseract and Image:
                            text = ocr_image_bytes(image)
                            if text:
                                chunks.append(text)
                except Exception:
                    pass
        return "\n".join(chunk.strip() for chunk in chunks if chunk and chunk.strip())
    except Exception:
        return ""


def read_html_bytes(data: bytes) -> str:
    try:
        try:
            html = data.decode("utf-8")
        except Exception:
            html = data.decode(errors="ignore")
        soup = BeautifulSoup(html, "html.parser")
        for tag in soup(["script", "style", "noscript", "template", "svg", "canvas", "math", "iframe", "header", "footer", "nav", "aside"]):
            tag.decompose()
        title = (soup.title.string or "").strip() if soup.title else ""
        main = soup.find(["main", "article"]) or soup.body or soup
        parts: List[str] = []
        if title:
            parts.append(title)

        def push(value: str) -> None:
            cleaned = re.sub(r"\s+", " ", (value or "").strip())
            if cleaned:
                parts.append(cleaned)

        for element in main.find_all(["h1", "h2", "h3", "h4", "h5", "h6", "p", "li", "figcaption"]):
            for image in element.find_all("img"):
                alt = (image.get("alt") or "").strip()
                if alt:
                    push(alt)
            push(element.get_text(separator=" ", strip=True))

        if not parts:
            push(main.get_text(separator=" ", strip=True))
        out = "\n".join(parts)
        return re.sub(r"\n{3,}", "\n\n", out).strip()
    except Exception:
        return ""


def ocr_image_bytes(data: bytes) -> str:
    if not pytesseract or not Image:
        return ""
    try:
        img = Image.open(io.BytesIO(data))
        return pytesseract.image_to_string(img)
    except Exception:
        return ""


def _normalize_ocr_text(text: str) -> str:
    return re.sub(r"[ \t\x0b\f\r]+", " ", (text or "")).strip()


def _guess_image_mime(filename: str, data: bytes) -> str:
    guessed, _ = mimetypes.guess_type(filename or "")
    if guessed and guessed.startswith("image/"):
        return guessed
    if data.startswith(b"\x89PNG\r\n\x1a\n"):
        return "image/png"
    if data.startswith(b"\xff\xd8"):
        return "image/jpeg"
    if data.startswith((b"GIF87a", b"GIF89a")):
        return "image/gif"
    if data[:4] == b"RIFF" and data[8:12] == b"WEBP":
        return "image/webp"
    if data.startswith((b"II*\x00", b"MM\x00*")):
        return "image/tiff"
    if data.startswith(b"BM"):
        return "image/bmp"
    return "image/png"


def _normalize_image_for_vlm(data: bytes, filename: str) -> tuple[bytes, str]:
    mime = _guess_image_mime(filename, data)
    if mime in {"image/png", "image/jpeg", "image/gif", "image/webp"}:
        return data, mime
    if not Image:
        return data, mime
    try:
        image = Image.open(io.BytesIO(data))
        if image.mode not in {"RGB", "L"}:
            image = image.convert("RGB")
        buffer = io.BytesIO()
        image.save(buffer, format="PNG")
        return buffer.getvalue(), "image/png"
    except Exception:
        return data, mime


def _to_data_url(data: bytes, mime: str) -> str:
    return f"data:{mime};base64,{base64.b64encode(data).decode('ascii')}"


def _build_vlm_payload(label: str, data: bytes, filename: str) -> Optional[dict[str, str]]:
    if not data:
        return None
    normalized, mime = _normalize_image_for_vlm(data, filename)
    return {
        "label": label,
        "url": _to_data_url(normalized, mime),
        "detail": "auto",
    }


def _combine_sections(sections: List[tuple[str, str]]) -> str:
    blocks = []
    for title, value in sections:
        cleaned = (value or "").strip()
        if cleaned:
            blocks.append(f"[{title}]\n{cleaned}")
    return "\n\n".join(blocks).strip()


def _render_pdf_pages_with_fitz(data: bytes, max_pages: int) -> List[dict[str, str]]:
    if fitz is None:
        return []
    payloads: List[dict[str, str]] = []
    doc = None
    try:
        doc = fitz.open(stream=data, filetype="pdf")
        page_count = min(int(getattr(doc, "page_count", 0) or 0), max_pages)
        for index in range(page_count):
            page = doc.load_page(index)
            pixmap = page.get_pixmap(matrix=fitz.Matrix(1.4, 1.4), alpha=False)
            payload = _build_vlm_payload(
                f"page {index + 1}",
                pixmap.tobytes("png"),
                f"page-{index + 1}.png",
            )
            if payload is not None:
                payloads.append(payload)
    except Exception:
        return []
    finally:
        if doc is not None:
            try:
                doc.close()
            except Exception:
                pass
    return payloads


def _extract_pdf_embedded_images(data: bytes, max_pages: int) -> List[dict[str, str]]:
    payloads: List[dict[str, str]] = []
    try:
        reader = PdfReader(io.BytesIO(data))
        for page_index, page in enumerate(reader.pages):
            if page_index >= max_pages:
                break
            images = list(getattr(page, "images", []) or [])
            if not images:
                continue
            for image_index, image in enumerate(images):
                raw = getattr(image, "data", None)
                if not isinstance(raw, bytes) or not raw:
                    pil_image = getattr(image, "image", None)
                    if pil_image is not None:
                        try:
                            buffer = io.BytesIO()
                            pil_image.save(buffer, format="PNG")
                            raw = buffer.getvalue()
                        except Exception:
                            raw = None
                if isinstance(raw, bytes) and raw:
                    payload = _build_vlm_payload(
                        f"page {page_index + 1}",
                        raw,
                        getattr(image, "name", None) or f"page-{page_index + 1}-image-{image_index + 1}.png",
                    )
                    if payload is not None:
                        payloads.append(payload)
                        break
    except Exception:
        return []
    return payloads[:max_pages]


def _extract_pdf_visual_payloads(data: bytes) -> List[dict[str, str]]:
    max_pages = max(
        1,
        min(
            int(getattr(settings, "VLM_MAX_IMAGES", 6)),
            int(getattr(settings, "VLM_PDF_PAGE_LIMIT", 4)),
        ),
    )
    payloads = _render_pdf_pages_with_fitz(data, max_pages)
    if payloads:
        return payloads[:max_pages]
    return _extract_pdf_embedded_images(data, max_pages)


def _summarize_image_bytes(filename: str, data: bytes, extracted_text: str = "") -> str:
    payload = _build_vlm_payload("image", data, filename)
    if payload is None:
        return ""
    return summarize_visual_material(filename, [payload], extracted_text=extracted_text)


def sniff_and_read(filename: str, data: bytes) -> str:
    lower = (filename or "").lower()
    if lower.endswith(".html") or lower.endswith(".htm"):
        return read_html_bytes(data)
    if lower.endswith(".pdf"):
        text = read_pdf_bytes(data)
        visual = ""
        if len(text.strip()) < 200:
            visual = summarize_visual_material(
                filename,
                _extract_pdf_visual_payloads(data),
                extracted_text=text,
            )
        if text and visual:
            return _combine_sections([
                ("Extracted text", text),
                ("Visual summary", visual),
            ])
        if text:
            return text
        if visual:
            return visual
        return _normalize_ocr_text(ocr_image_bytes(data))
    if lower.endswith(".docx"):
        return read_docx_bytes(data)
    if lower.endswith(".xlsx") or lower.endswith(".xls"):
        return read_xlsx_bytes(data)
    if lower.endswith(".pptx"):
        return read_pptx_bytes(data)
    if lower.endswith(".ppt"):
        try:
            return data.decode("utf-8", errors="ignore")
        except Exception:
            return ""
    if any(lower.endswith(ext) for ext in (".png", ".jpg", ".jpeg", ".bmp", ".tif", ".tiff", ".gif", ".webp")):
        visual = _summarize_image_bytes(filename, data)
        ocr = _normalize_ocr_text(ocr_image_bytes(data))
        if visual and ocr:
            return _combine_sections([
                ("Visual summary", visual),
                ("OCR fallback", ocr),
            ])
        return visual or ocr
    try:
        head = data[:512].decode("utf-8", errors="ignore").lower()
        if "<html" in head or "<!doctype html" in head:
            return read_html_bytes(data)
    except Exception:
        pass
    try:
        return data.decode("utf-8", errors="ignore")
    except Exception:
        return ""
